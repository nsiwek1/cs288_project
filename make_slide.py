"""Generate a single slide describing the Cycle Method for class presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# --- Background ---
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# --- Color palette ---
DARK = RGBColor(0x1B, 0x1B, 0x1B)
ACCENT = RGBColor(0x2E, 0x86, 0x4B)  # forest green
ACCENT2 = RGBColor(0xC0, 0x5C, 0x2E)  # warm orange
LIGHT_BG = RGBColor(0xF4, 0xF1, 0xEB)  # warm off-white
MUTED = RGBColor(0x66, 0x66, 0x66)
BLUE = RGBColor(0x2B, 0x6C, 0xB3)

def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.rotation = 0.0
    # Minimal corner rounding
    shape.adjustments[0] = 0.05
    return shape

# Helper conversions
def inch(x): return int(Inches(x))
def pt_to_emu(x): return int(Pt(x) * 12700 / 12700 * 914400 / 72)

# ============================================================
# TITLE BAR (green accent bar at top)
# ============================================================
bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    0, 0, inch(13.333), inch(0.9)
)
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

add_text(slide, inch(0.5), inch(0.12), inch(10), inch(0.7),
         "The Cycle Method: CycleGAN for Dialect Detection", size=30, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), font_name="Calibri")

# Subtitle
add_text(slide, inch(0.5), inch(1.0), inch(12), inch(0.45),
         "Using cycle-consistent neural networks to quantify acoustic differences in chimpanzee pant-hoots",
         size=15, color=MUTED, font_name="Calibri")

# ============================================================
# LEFT COLUMN — "What is CycleGAN?" + "Key Idea"
# ============================================================
col1_x = inch(0.4)

# Box 1: What is CycleGAN?
add_rounded_rect(slide, col1_x, inch(1.65), inch(4.0), inch(2.4), LIGHT_BG)
add_text(slide, col1_x + inch(0.2), inch(1.75), inch(3.6), inch(0.35),
         "What is CycleGAN?", size=16, bold=True, color=ACCENT)

what_box = slide.shapes.add_textbox(
    Emu(col1_x + inch(0.2)), Emu(inch(2.15)), Emu(inch(3.6)), Emu(inch(1.8))
)
tf = what_box.text_frame
tf.word_wrap = True

bullets = [
    "A neural network that learns to translate between two domains without paired examples",
    "Two generators learn opposite mappings:\n   G₁: Pant → Hoot    G₂: Hoot → Pant",
    "Enforces cycle consistency: translating back and forth should recover the original",
]
for i, b in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = b
    p.font.size = Pt(12)
    p.font.color.rgb = DARK
    p.font.name = "Calibri"
    p.space_after = Pt(6)

# Box 2: Key Idea
add_rounded_rect(slide, col1_x, inch(4.2), inch(4.0), inch(2.9), LIGHT_BG)
add_text(slide, col1_x + inch(0.2), inch(4.3), inch(3.6), inch(0.35),
         "Key Idea", size=16, bold=True, color=ACCENT)

idea_box = slide.shapes.add_textbox(
    Emu(col1_x + inch(0.2)), Emu(inch(4.7)), Emu(inch(3.6)), Emu(inch(2.3))
)
tf = idea_box.text_frame
tf.word_wrap = True

ideas = [
    "Train the cycle on one community (Kasekela)",
    "Apply it to calls from other communities",
    "Cycle reconstruction error = dialect distance",
    "Low error → fits learned pattern (same dialect)",
    "High error → doesn't fit (different dialect)",
]
for i, b in enumerate(ideas):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = "• " + b
    p.font.size = Pt(12)
    p.font.color.rgb = DARK
    p.font.name = "Calibri"
    p.space_after = Pt(5)

# Hypothesis highlight
hyp_box = slide.shapes.add_textbox(
    Emu(col1_x + inch(0.15)), Emu(inch(6.35)), Emu(inch(3.7)), Emu(inch(0.55))
)
tf = hyp_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Hypothesis: Cycle loss ordering"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = ACCENT2
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Kasekela < Mitumba < Kanyawara"
p2.font.size = Pt(14)
p2.font.bold = True
p2.font.color.rgb = ACCENT2
p2.alignment = PP_ALIGN.CENTER

# ============================================================
# CENTER COLUMN — Cycle Diagram
# ============================================================
col2_x = inch(4.7)

add_text(slide, col2_x + inch(0.3), inch(1.65), inch(4.0), inch(0.35),
         "The Cycle (Illustrated)", size=16, bold=True, color=ACCENT,
         alignment=PP_ALIGN.CENTER)

# Draw a simplified cycle diagram with shapes and arrows
# Pant box
pant_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Emu(col2_x + inch(0.3)), Emu(inch(2.6)), Emu(inch(1.5)), Emu(inch(0.7))
)
pant_shape.fill.solid()
pant_shape.fill.fore_color.rgb = ACCENT
pant_shape.line.fill.background()
pant_tf = pant_shape.text_frame
pant_tf.paragraphs[0].text = "Pant\n(Buildup)"
pant_tf.paragraphs[0].font.size = Pt(13)
pant_tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
pant_tf.paragraphs[0].font.bold = True
pant_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
pant_tf.word_wrap = True

# Hoot box
hoot_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Emu(col2_x + inch(2.8)), Emu(inch(2.6)), Emu(inch(1.5)), Emu(inch(0.7))
)
hoot_shape.fill.solid()
hoot_shape.fill.fore_color.rgb = ACCENT2
hoot_shape.line.fill.background()
hoot_tf = hoot_shape.text_frame
hoot_tf.paragraphs[0].text = "Hoot\n(Climax)"
hoot_tf.paragraphs[0].font.size = Pt(13)
hoot_tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
hoot_tf.paragraphs[0].font.bold = True
hoot_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
hoot_tf.word_wrap = True

# Reconstructed Pant box
recon_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Emu(col2_x + inch(1.55)), Emu(inch(4.1)), Emu(inch(1.5)), Emu(inch(0.7))
)
recon_shape.fill.solid()
recon_shape.fill.fore_color.rgb = ACCENT
recon_shape.line.fill.background()
recon_shape.line.color.rgb = ACCENT
recon_tf = recon_shape.text_frame
recon_tf.paragraphs[0].text = "Reconstructed\nPant ≈ Original?"
recon_tf.paragraphs[0].font.size = Pt(11)
recon_tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
recon_tf.paragraphs[0].font.bold = True
recon_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
recon_tf.word_wrap = True

# Arrow labels
add_text(slide, col2_x + inch(1.6), inch(2.2), inch(1.4), inch(0.35),
         "G₁: Pant → Hoot", size=11, bold=True, color=BLUE,
         alignment=PP_ALIGN.CENTER)

add_text(slide, col2_x + inch(0.15), inch(3.55), inch(1.6), inch(0.35),
         "G₂: Hoot → Pant", size=11, bold=True, color=BLUE,
         alignment=PP_ALIGN.CENTER)

add_text(slide, col2_x + inch(2.85), inch(3.55), inch(1.6), inch(0.35),
         "G₂: Hoot → Pant", size=11, bold=True, color=BLUE,
         alignment=PP_ALIGN.CENTER)

# Arrow connectors (using simple lines)
# Top arrow: Pant -> Hoot
arrow1 = slide.shapes.add_connector(
    1,  # straight
    Emu(col2_x + inch(1.8)), Emu(inch(2.95)),
    Emu(col2_x + inch(2.8)), Emu(inch(2.95))
)
arrow1.line.color.rgb = BLUE
arrow1.line.width = Pt(2)

# Right arrow: Hoot -> Reconstructed
arrow2 = slide.shapes.add_connector(
    1,
    Emu(col2_x + inch(3.55)), Emu(inch(3.3)),
    Emu(col2_x + inch(3.05)), Emu(inch(4.1))
)
arrow2.line.color.rgb = BLUE
arrow2.line.width = Pt(2)

# Left arrow: also shows flow
arrow3 = slide.shapes.add_connector(
    1,
    Emu(col2_x + inch(1.05)), Emu(inch(3.3)),
    Emu(col2_x + inch(1.55)), Emu(inch(4.1))
)
arrow3.line.color.rgb = BLUE
arrow3.line.width = Pt(2)

# Cycle loss annotation
loss_box = slide.shapes.add_textbox(
    Emu(col2_x + inch(0.5)), Emu(inch(5.1)), Emu(inch(3.6)), Emu(inch(0.9))
)
tf = loss_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Cycle Loss = | Original Pant − Reconstructed Pant |"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Low loss = good fit  •  High loss = different dialect"
p2.font.size = Pt(11)
p2.font.color.rgb = MUTED
p2.alignment = PP_ALIGN.CENTER

# Embedding note
add_text(slide, col2_x + inch(0.3), inch(6.15), inch(4.0), inch(0.7),
         "Operates on 768-dim AVEX embeddings extracted from audio recordings, not on raw audio",
         size=10, color=MUTED, alignment=PP_ALIGN.CENTER, font_name="Calibri")

# ============================================================
# RIGHT COLUMN — Pipeline Steps + Stats
# ============================================================
col3_x = inch(9.2)

add_rounded_rect(slide, col3_x, inch(1.65), inch(3.8), inch(3.0), LIGHT_BG)
add_text(slide, col3_x + inch(0.2), inch(1.75), inch(3.4), inch(0.35),
         "Pipeline", size=16, bold=True, color=ACCENT)

pipe_box = slide.shapes.add_textbox(
    Emu(col3_x + inch(0.2)), Emu(inch(2.15)), Emu(inch(3.4)), Emu(inch(2.4))
)
tf = pipe_box.text_frame
tf.word_wrap = True

steps = [
    ("1.", "Extract AVEX embeddings from .wav files"),
    ("2.", "Standardize embeddings per community"),
    ("3.", "Train CycleGAN on Kasekela (5-fold CV)"),
    ("4.", "Apply to Kasekela, Mitumba, Kanyawara"),
    ("5.", "Compare cycle reconstruction loss"),
    ("6.", "Statistical tests (Mann-Whitney U, permutation tests, bootstrap CI)"),
]
for i, (num, step) in enumerate(steps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{num}  {step}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK
    p.font.name = "Calibri"
    p.space_after = Pt(5)

# Training details box
add_rounded_rect(slide, col3_x, inch(4.8), inch(3.8), inch(2.3), LIGHT_BG)
add_text(slide, col3_x + inch(0.2), inch(4.9), inch(3.4), inch(0.35),
         "Model Details", size=16, bold=True, color=ACCENT)

detail_box = slide.shapes.add_textbox(
    Emu(col3_x + inch(0.2)), Emu(inch(5.3)), Emu(inch(3.4)), Emu(inch(1.7))
)
tf = detail_box.text_frame
tf.word_wrap = True

details = [
    "Generators: 4-layer MLP (512 hidden dims)",
    "Cycle weight: λ = 20.0 (heavy emphasis)",
    "5-fold cross-validation, early stopping",
    "Unpaired learning — no matched pant-hoot pairs needed",
    "Validated via within-community sanity checks (SL vs. GIM opposing dialects)",
]
for i, d in enumerate(details):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = "• " + d
    p.font.size = Pt(11)
    p.font.color.rgb = DARK
    p.font.name = "Calibri"
    p.space_after = Pt(4)

# ============================================================
# Save
# ============================================================
out_path = "/Users/natalia_mac/Downloads/Chimpanzee Pant-Hoots/Cycle_Method_Slide.pptx"
prs.save(out_path)
print(f"Saved to: {out_path}")
